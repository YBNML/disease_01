"""Generate a project presentation (.pptx) summarizing disease_01.

Usage:
    python scripts/make_presentation.py
    # → docs/disease_01_presentation.pptx
"""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = PROJECT_ROOT / "docs" / "disease_01_presentation.pptx"
ASSETS = PROJECT_ROOT / "docs" / "ppt_assets"

# Colors
COLOR_PRIMARY = RGBColor(0x16, 0x5B, 0x33)   # dark green (citrus leaf)
COLOR_ACCENT = RGBColor(0xE8, 0x72, 0x00)    # orange (citrus fruit)
COLOR_TEXT = RGBColor(0x1F, 0x1F, 0x1F)
COLOR_MUTED = RGBColor(0x55, 0x55, 0x55)
COLOR_BG = RGBColor(0xFA, 0xFA, 0xFA)
COLOR_HEADER_BG = RGBColor(0x16, 0x5B, 0x33)
COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)

# Font
FONT = "Malgun Gothic"   # Korean-capable; PowerPoint will substitute on macOS


def _set_font(run, size_pt: int, bold: bool = False, color: RGBColor = COLOR_TEXT):
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background rectangle (top band)
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(2.4))
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), prs.slide_width - Inches(1.4), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "감귤 궤양병 CV 파이프라인"
    _set_font(r, 40, bold=True, color=COLOR_HEADER_FG)

    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "disease_01 — Classification / Detection / Segmentation + 비교 실험"
    _set_font(r, 20, color=COLOR_HEADER_FG)

    # Subtitle area
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), prs.slide_width - Inches(1.4), Inches(3.2))
    tf = tx.text_frame
    tf.word_wrap = True

    lines = [
        ("YBNML", 18, True, COLOR_PRIMARY),
        ("AI Hub 감귤 병충해 데이터 / Apple M4 MPS 환경", 14, False, COLOR_TEXT),
        ("", 8, False, COLOR_TEXT),
        ("프로젝트 목표 · 알고리즘 · 결과 · 인사이트", 16, True, COLOR_TEXT),
        ("GitHub: github.com/YBNML/disease_01", 13, False, COLOR_MUTED),
        ("", 8, False, COLOR_TEXT),
        ("발표 일시: 2026-04-21", 12, False, COLOR_MUTED),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = text
        _set_font(r, size, bold=bold, color=color)
        p.space_after = Pt(4)


def add_section_header(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), prs.slide_width - Inches(1.4), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 40, bold=True, color=COLOR_HEADER_FG)
    if subtitle:
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = subtitle
        _set_font(r, 18, color=COLOR_HEADER_FG)


def add_content_slide(prs, title: str, bullets: list, notes: str = ""):
    """bullets: list of (text, level, bold?) or just str."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 28, bold=True, color=COLOR_PRIMARY)

    # Underline stripe
    from pptx.enum.shapes import MSO_SHAPE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Emu(38100))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_ACCENT
    stripe.line.fill.background()

    # Content
    body_top = Inches(1.4)
    body = slide.shapes.add_textbox(Inches(0.5), body_top, prs.slide_width - Inches(1.0),
                                     prs.slide_height - body_top - Inches(0.3))
    tf = body.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        if isinstance(item, str):
            text, level, bold = item, 0, False
        else:
            text, level, bold = (item + (False,))[:3] if len(item) == 2 else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        bullet_char = "•" if level == 0 else "–"
        indent = "    " * level
        r = p.add_run()
        r.text = f"{indent}{bullet_char}  {text}" if text else ""
        size = 18 if level == 0 else 15
        color = COLOR_TEXT if level == 0 else COLOR_MUTED
        _set_font(r, size, bold=bold, color=color)
        p.space_after = Pt(6)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _add_title_stripe(slide, prs, title: str):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 26, bold=True, color=COLOR_PRIMARY)
    from pptx.enum.shapes import MSO_SHAPE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.5), Emu(38100))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_ACCENT
    stripe.line.fill.background()


def _caption(slide, x, y, w, text, size=12, italic=True, align=PP_ALIGN.CENTER):
    tx = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    _set_font(r, size, color=COLOR_MUTED)
    r.font.italic = italic


def add_single_image_slide(prs, title: str, image_path, caption: str = "", notes: str = ""):
    """One large centered image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_stripe(slide, prs, title)

    # Max box
    max_w = prs.slide_width - Inches(1.0)
    max_h = prs.slide_height - Inches(2.0)
    pic = slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.4), width=max_w)
    # If image too tall, scale down
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = int(pic.height * ratio)
    # Center horizontally
    pic.left = int((prs.slide_width - pic.width) / 2)

    if caption:
        _caption(slide, Inches(0.5), pic.top + pic.height + Inches(0.1),
                 prs.slide_width - Inches(1.0), caption)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_two_image_slide(prs, title: str, left_path, right_path,
                        left_caption: str = "", right_caption: str = "",
                        notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_stripe(slide, prs, title)

    gap = Inches(0.3)
    col_w = (prs.slide_width - Inches(1.0) - gap) / 2
    max_h = prs.slide_height - Inches(2.2)
    top = Inches(1.4)

    for i, (path, cap) in enumerate([(left_path, left_caption), (right_path, right_caption)]):
        left = Inches(0.5) + i * (col_w + gap)
        pic = slide.shapes.add_picture(str(path), left, top, width=int(col_w))
        if pic.height > max_h:
            ratio = max_h / pic.height
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
            # Recenter horizontally inside the column
            pic.left = int(left + (col_w - pic.width) / 2)
        if cap:
            _caption(slide, left, top + pic.height + Inches(0.1), int(col_w), cap, size=11)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_table_slide(prs, title: str, header: list, rows: list, notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set_font(r, 28, bold=True, color=COLOR_PRIMARY)

    from pptx.enum.shapes import MSO_SHAPE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(1.5), Emu(38100))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLOR_ACCENT
    stripe.line.fill.background()

    n_rows = len(rows) + 1
    n_cols = len(header)
    left, top = Inches(0.5), Inches(1.5)
    width = prs.slide_width - Inches(1.0)
    height = Inches(5.2)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header row
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = ""
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(h)
        _set_font(r, 13, bold=True, color=COLOR_HEADER_FG)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            # right-align numbers (simple heuristic)
            s = str(val)
            is_number = s.replace(",", "").replace(".", "").replace("%", "").replace("-", "").replace("+", "").replace(" ", "").isdigit()
            p.alignment = PP_ALIGN.RIGHT if j > 0 and is_number else PP_ALIGN.LEFT
            r = p.add_run(); r.text = s
            _set_font(r, 12, color=COLOR_TEXT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG if i % 2 == 1 else RGBColor(0xEE, 0xEE, 0xEE)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 1. Title ---
    add_title_slide(prs)

    # --- 2. Agenda ---
    add_content_slide(prs, "목차",
        [
            "1. 프로젝트 개요 (목적 · 데이터)",
            "2. 아키텍처 & 파이프라인 설계",
            "3. 3가지 CV 태스크 구현",
            ("P1 Classification / P2 Detection / P3 Segmentation", 1),
            "4. 비교 실험 3종",
            ("P1b 백본 / P2b YOLO / P3b Seg 아키텍처", 1),
            "5. 배포 실험 (장기학습 · 양자화 · Distillation)",
            "6. 분석 노트북 & 핵심 인사이트",
            "7. 한계 & 향후 계획",
        ])

    # --- 3. 프로젝트 개요 ---
    add_section_header(prs, "1. 프로젝트 개요", "목적과 데이터셋")

    add_content_slide(prs, "프로젝트 목적",
        [
            ("문제 정의", 0, True),
            ("감귤(온주밀감) 이미지로 정상 vs 궤양병(canker) 자동 감별", 1),
            ("필요성", 0, True),
            ("궤양병은 수확량·상품성에 직접 타격", 1),
            ("조기 감별로 선별·검역 비용 절감", 1),
            ("접근 방법", 0, True),
            ("3가지 CV 태스크를 공통 모듈 위에 구축", 1),
            ("모델 아키텍처 비교로 배포 관점 평가", 1),
            ("TDD 기반 재현 가능한 실험 프레임워크", 1),
        ])

    add_content_slide(prs, "데이터셋: AI Hub 감귤 병충해 영상",
        [
            "출처: AI Hub 감귤 병충해 영상 데이터 (온주밀감)",
            "총 3,834장 — 정상 2,290 / 궤양병 1,544",
            ("Train 3,407 / Val 427 (AI Hub 제공 split 유지)", 1),
            ("Polygon 라벨 보유: 787장 (Detection/Segmentation용)", 1),
            "이미지 특성: 1920×1080, 배경 흰색 사전 처리됨",
            ("이미지당 감귤 1개 (단일 객체)", 1),
            "풍부한 메타데이터",
            ("카메라(Samsung/Xiaomi/LGE), 지역(10곳), 노지/온실", 1),
            ("환경(기온/습도/일사량/토양수분/강수)", 1),
            ("성장단계, 촬영시기", 1),
        ])

    # --- 4. Architecture ---
    add_section_header(prs, "2. 아키텍처", "설계 원칙과 모듈 구성")

    add_content_slide(prs, "설계 원칙",
        [
            ("DRY — 공통 모듈 재사용", 0, True),
            ("데이터 로딩 · 라벨 파싱 · 설정 · 유틸을 common/ 하나로 통일", 1),
            ("태스크별 독립성", 0, True),
            ("classification/ detection/ segmentation/ 각 태스크가 common/만 의존", 1),
            ("TDD — 총 93개 테스트", 0, True),
            ("라벨 파싱 / Dataset / Metrics / Loss 모두 단위 테스트", 1),
            ("재현성", 0, True),
            ("고정 seed(42) + 타임스탬프 output + config 스냅샷", 1),
            ("Config-driven 실험", 0, True),
            ("YAML + CLI override로 하이퍼파라미터 실험", 1),
        ])

    add_content_slide(prs, "데이터 파이프라인 & 라벨 변환",
        [
            ("AI Hub JSON → Python dict (load_sample)", 0, True),
            "ANTN_PT 문자열 파싱: '[x1|x2|...],[y1|y2|...]' → (N, 2) int array",
            ("태스크별 변환", 0, True),
            ("Classification: 폴더명 → 0(정상) / 1(궤양병)", 1),
            ("Detection: polygon → bbox → YOLO format (정규화)", 1),
            ("Segmentation: polygon → 픽셀 mask (0:bg, 1:normal, 2:canker)", 1),
            ("메타데이터 동시 반환 (Phase 2 대비)", 0, True),
            ("카메라 · 지역 · 환경값을 Dataset 샘플에 포함", 1),
        ])

    # --- 5. 기본 3태스크 ---
    add_section_header(prs, "3. 3가지 CV 태스크 구현", "P1 Classification / P2 Detection / P3 Segmentation")

    add_single_image_slide(
        prs, "결과 이미지 — ViT-Small Confusion Matrix",
        ASSETS / "cls_vit_cm.png",
        caption="ViT-Small/16 on val set (427) — 오분류 4장 (FN 1, FP 3)",
        notes="P1b에서 최고 성능을 보인 ViT-Small의 confusion matrix. 궤양병 recall 거의 완벽.",
    )

    add_content_slide(prs, "P1 — Classification (ResNet50)",
        [
            ("모델 & 학습 설정", 0, True),
            ("torchvision ResNet50 (ImageNet pretrained), FC → 2-class", 1),
            ("AdamW (lr=1e-4) + CosineAnnealingLR + CrossEntropyLoss", 1),
            ("WeightedRandomSampler로 클래스 불균형 처리", 1),
            ("Best checkpoint: val F1 (궤양병 class)", 1),
            ("결과 — 2 epochs on MPS (약 6분)", 0, True),
            ("Accuracy 98.83% / F1 0.986 / AUC 0.999", 1),
            ("궤양병 Recall 100% (놓친 사례 0장)", 1),
            ("FP 5장 / FN 0장 — 의료적 스크리닝에 이상적", 1),
        ])

    add_two_image_slide(
        prs, "결과 이미지 — YOLOv8m 30ep 학습 곡선 & Confusion Matrix",
        ASSETS / "det_yolov8m_training.png",
        ASSETS / "det_yolov8m_cm.png",
        left_caption="Training curves (loss & mAP over 30 epochs)",
        right_caption="Confusion matrix (val 88 images)",
        notes="P4-A1 장기학습 결과: under-fit이었던 yolov8m이 30ep에 mAP@0.5 0.9945 달성.",
    )

    add_content_slide(prs, "P2 — Detection (YOLOv8s)",
        [
            ("모델 & 학습 설정", 0, True),
            ("ultralytics YOLOv8s (COCO pretrained)", 1),
            ("AI Hub polygon → YOLO bbox 포맷 변환 (787장)", 1),
            ("SGD + 내장 augmentation (mosaic/mixup)", 1),
            ("imgsz=640, batch=16", 1),
            ("결과 — 5 epochs on MPS (약 6분)", 0, True),
            ("mAP@0.5: 0.994 / mAP@0.5:0.95: 0.989", 1),
            ("Normal AP 0.994 / Canker AP 0.995", 1),
            ("배경 흰색 + 단일 객체 → 난이도 본질적으로 낮음", 1),
        ])

    add_two_image_slide(
        prs, "결과 이미지 — Segmentation 정성 샘플 (DeepLabV3+ 20ep)",
        ASSETS / "seg_sample_000.png",
        ASSETS / "seg_sample_001.png",
        left_caption="샘플 #1 (원본 | GT mask | 예측 mask)",
        right_caption="샘플 #2 (원본 | GT mask | 예측 mask)",
        notes="P4-A2 결과. 초록=정상, 빨강=궤양병. GT와 예측이 거의 겹침.",
    )

    add_content_slide(prs, "P3 — Segmentation (smp U-Net + ResNet34)",
        [
            ("모델 & 학습 설정", 0, True),
            ("smp U-Net encoder=ResNet34 (ImageNet pretrained)", 1),
            ("3-class semantic: bg / 정상 / 궤양병 — 단일 모델로 분류+분할", 1),
            ("AdamW + Cosine + 0.5×CE + 0.5×DiceLoss", 1),
            ("albumentations로 image+mask 동시 변환 (mask NN interp)", 1),
            ("결과 — 3 epochs on MPS (약 8분)", 0, True),
            ("mIoU 0.940 / Pixel Accuracy 99.3%", 1),
            ("IoU: bg 0.998 / 정상 0.937 / 궤양병 0.886", 1),
            ("궤양병 최약 — polygon이 병변 아닌 과일 외곽이라는 한계", 1),
        ])

    # --- 6. 비교 실험 3종 ---
    add_section_header(prs, "4. 비교 실험 3종", "백본 · YOLO variants · Seg 아키텍처")

    add_single_image_slide(
        prs, "결과 이미지 — P1b 백본 비교 시각화",
        ASSETS / "nb05_fig_02.png",
        caption="scatter: params vs accuracy / latency vs F1 / throughput vs accuracy",
        notes="노트북 05에서 실행한 3-panel scatter. ViT-Small이 Pareto-optimal.",
    )

    add_table_slide(prs, "P1b — Classification 백본 비교 (5 models, 5 epochs)",
        ["Model", "Params", "Accuracy", "F1 (canker)", "Latency bs=1 (ms)", "FPS"],
        [
            ["ViT-Small/16 🏆", "21.7M", "99.06%", "0.988", "5.94", "175.0"],
            ["ConvNeXt-Tiny", "27.8M", "98.83%", "0.986", "7.18", "109.8"],
            ["ResNet50", "23.5M", "98.83%", "0.986", "27.86", "83.8"],
            ["MobileNetV3-L", "4.2M", "97.66%", "0.971", "13.12", "264.6"],
            ["EfficientNet-B0", "4.0M", "97.42%", "0.968", "19.68", "160.6"],
        ],
        notes="ViT-Small이 정확도 최고 + latency 최저 — Apple MPS에서 Transformer가 기존 CNN보다 빠르게 동작한 흥미로운 발견. MobileNetV3은 throughput 왕 (264 FPS, 4.2M params).",
    )

    add_table_slide(prs, "P2b — YOLOv8 variants 비교 (n/s/m, 5 epochs)",
        ["Model", "Params", "mAP@0.5", "mAP@0.5:0.95", "Latency (ms)", "FPS"],
        [
            ["yolov8n", "3.0M", "0.9926", "0.9907", "8.87", "168.9"],
            ["yolov8s 🏆", "11.1M", "0.9933", "0.9917", "17.24", "68.3"],
            ["yolov8m (5ep under-fit)", "25.9M", "0.9909", "0.9860", "34.91", "30.2"],
        ],
        notes="5 epoch로는 yolov8m이 under-fit. yolov8s가 정확도+속도 균형 최고. yolov8n은 엣지/모바일 유력 (168 FPS).",
    )

    add_table_slide(prs, "P3b — Segmentation 아키텍처 비교 (4 combos, 3 epochs)",
        ["Label", "Arch", "Encoder", "Params", "mIoU", "IoU canker", "FPS"],
        [
            ["unet_resnet34 (baseline)", "Unet", "resnet34", "24.4M", "0.9616", "0.9291", "40.5"],
            ["deeplabv3plus_resnet34 🏆", "DeepLabV3+", "resnet34", "22.4M", "0.9777", "0.9621", "36.1"],
            ["fpn_resnet34", "FPN", "resnet34", "23.2M", "0.9677", "0.9400", "41.6"],
            ["unet_efficientnet-b0", "Unet", "efficientnet-b0", "6.3M", "0.9077", "0.8172", "43.3"],
        ],
        notes="DeepLabV3+ 우승. Atrous conv로 공간 정보 복원이 병변 경계에 유리. Decoder 교체 > Encoder 경량화.",
    )

    # --- 7. 배포 실험 (P4) ---
    add_section_header(prs, "5. 배포 실험 (P4)", "장기학습 · 양자화 · Knowledge Distillation")

    add_content_slide(prs, "P4-A1 — yolov8m 30 epoch (under-fit 가설 검증)",
        [
            ("P2b에서의 의문", 0, True),
            ("5ep에 yolov8m(0.991)이 yolov8s(0.993)보다 낮음 — under-fit?", 1),
            ("검증: 30 epoch로 재학습", 0, True),
            ("Epoch 5 → 10 → 20 → 30 mAP@0.5: 0.927 → 0.991 → 0.988 → 0.9945", 1),
            ("결론", 0, True),
            ("Under-fit 가설 확인 — 충분한 학습이 있어야 모델 capacity 발휘", 1),
            ("큰 모델은 '긴 학습'이 조건", 1),
        ])

    add_content_slide(prs, "P4-A2 — DeepLabV3+ 20 epoch (장기 학습)",
        [
            ("P3b 3ep vs P4 20ep 비교", 0, True),
            ("mIoU: 0.9777 → 0.9836 (+0.006)", 1),
            ("궤양병 IoU: 0.9621 → 0.9721 (+0.010)", 1),
            ("정상 IoU: 0.9732 → 0.9805 (+0.007)", 1),
            ("17~20 epoch 구간 수렴 관찰", 0, True),
            ("추가 개선을 위해서는 데이터·라벨 품질 개선이 필요", 1),
            ("교훈", 0, True),
            ("장기 학습의 ROI는 데이터 난이도에 강하게 의존", 1),
        ])

    add_table_slide(prs, "P4-G — ViT-Small INT8 양자화",
        ["Variant", "Size (MB)", "Accuracy", "F1 (canker)", "Latency bs=1 (ms)", "Throughput (FPS)"],
        [
            ["fp32 (baseline)", "82.7", "0.9906", "0.9884", "18.5", "88.9"],
            ["int8 (dynamic)", "22.0 (−73%)", "0.9906 (동일)", "0.9884 (동일)", "32.2 (+74%)", "38.5 (−57%)"],
        ],
        notes="크기 73% 축소 + 정확도 완전 보존. 그러나 CPU에서 dynamic quant 오버헤드로 latency 오히려 증가. 교훈: 양자화 ≠ 항상 빠름. 크기 제약 환경에서 유용.",
    )

    add_content_slide(prs, "P4-H — Knowledge Distillation (ViT → MobileNetV3)",
        [
            ("구성", 0, True),
            ("Teacher: ViT-Small/16 (21.7M, F1=0.988) — frozen", 1),
            ("Student: MobileNetV3-Large (4.2M)", 1),
            ("Loss: α×CE + (1-α)×T²×KL, α=0.5, T=4.0", 1),
            ("결과 (5 epoch)", 0, True),
            ("Student F1 = 0.962  (plain MobileNetV3 P1b는 0.971)", 1),
            ("Distillation이 plain보다 오히려 약간 낮음", 1),
            ("해석", 0, True),
            ("α/T 하이퍼파라미터 튜닝 미흡", 1),
            ("데이터가 쉬워 학생이 이미 충분 → KD 이득 제한", 1),
            ("교훈: KD는 마법이 아니다. 적절한 세팅·데이터 난이도 필요", 1),
        ])

    # --- 8. 분석 노트북 & 인사이트 ---
    add_section_header(prs, "6. 분석 & 인사이트", "Grad-CAM / FP 분석 / 핵심 발견")

    add_single_image_slide(
        prs, "결과 이미지 — Grad-CAM (정상 샘플)",
        ASSETS / "nb07_gradcam_01.png",
        caption="4장의 정상 감귤 샘플에 Grad-CAM overlay — 모델이 어디를 보는가",
        notes="ResNet50이 정상 샘플에서 과일 중앙~표면 텍스처에 집중하는 패턴.",
    )

    add_single_image_slide(
        prs, "결과 이미지 — Grad-CAM (궤양병 샘플)",
        ASSETS / "nb07_gradcam_02.png",
        caption="궤양병 감귤 샘플 Grad-CAM — 병변 추정 영역에 activation 집중",
        notes="병변성 패턴을 가진 영역에 모델 attention이 몰림.",
    )

    add_content_slide(prs, "분석 노트북 8종 (docs/analysis/)",
        [
            "01 Dataset EDA — 클래스 분포 · polygon 커버리지 · 메타데이터",
            "02 Classification 심층 — 학습 곡선 · CM · 오분류 · ROC",
            "03 Detection 심층 — bbox 시각화 · per-class AP · 실패 케이스",
            "04 Segmentation 심층 — mask 시각화 · IoU/Dice · failure cases",
            "05 Backbone 비교 심층 — Pareto frontier · 배포 시나리오 (실행 완료)",
            "06 Lessons learned — 설계 근거 · 트러블슈팅 · 다음 계획",
            "07 Grad-CAM 시각화 — 모델이 어디를 보고 판단?",
            "08 FP 오분류 분석 — 4장 FP 공통점 탐색",
        ])

    add_content_slide(prs, "FP(False Positive) 오분류 분석 결과",
        [
            ("ResNet50 val set에서 오분류 5장 (정상→궤양병 4 + 궤양병→정상 1)", 0, True),
            ("4장 FP 공통점", 0, True),
            ("모두 Samsung(3) 또는 LGE(1) 카메라 — Xiaomi 없음", 1),
            ("고습 조건 (습도 69~100%)", 1),
            ("예측 확률 0.94+ — 고신뢰도 오판 (hard cases)", 1),
            ("일부는 저조도 + 늦가을 저온 조건", 1),
            ("시사점", 0, True),
            ("환경·카메라 condition이 모델의 실제 에러 패턴을 설명 가능", 1),
            ("향후 멀티모달 / 도메인 분석에 실마리 제공", 1),
        ])

    add_content_slide(prs, "핵심 인사이트 (1/2)",
        [
            ("1. 데이터 난이도가 실험 설계를 규정", 0, True),
            ("모든 모델이 97%+ → 정확도 차별화는 미미", 1),
            ("실질 차별화는 속도 · 크기에서 발생", 1),
            ("2. Pretraining의 힘", 0, True),
            ("ImageNet pretrained만으로 2 epoch에 98% 이상 수렴", 1),
            ("Scratch 학습 대비 시간·성능 둘 다 압도", 1),
            ("3. Transformer의 MPS 효율성", 0, True),
            ("ViT-Small이 conv 모델(ResNet/ConvNeXt)보다 빠름", 1),
            ("Apple Neural Engine이 attention 연산에 유리한 것으로 추정", 1),
        ])

    add_content_slide(prs, "핵심 인사이트 (2/2)",
        [
            ("4. 장기 학습의 가치", 0, True),
            ("큰 모델일수록 under-fit 가능성 ↑", 1),
            ("yolov8m 5ep → 30ep로 mAP +0.003 상승", 1),
            ("5. 양자화 ≠ 항상 빠름", 0, True),
            ("Apple Silicon + 작은 ViT에서 dynamic INT8 오히려 느림", 1),
            ("크기 이득은 분명하나 속도는 컨텍스트 의존적", 1),
            ("6. Distillation의 한계", 0, True),
            ("학생이 이미 충분히 좋으면 KD 이득 제한", 1),
            ("α/T 하이퍼파라미터 튜닝이 결정적", 1),
            ("7. 라벨 한계의 본질적 제약", 0, True),
            ("Polygon이 과일 외곽선 — 진짜 병변 위치가 아님", 1),
            ("궤양병 IoU가 모든 실험에서 최약 → 라벨 품질이 천장", 1),
        ])

    # --- 9. 한계 / Future ---
    add_section_header(prs, "7. 한계 & 향후 계획", "Limitations and Future Work")

    add_content_slide(prs, "한계 (Limitations)",
        [
            ("라벨의 의미적 한계", 0, True),
            ("Polygon은 과일 외곽선이지 병변 영역이 아님", 1),
            ("진정한 '병변 분할'은 별도 라벨링 필요", 1),
            ("데이터가 쉬움", 0, True),
            ("배경 흰색 사전처리 — 자연 배경 robustness 미검증", 1),
            ("이미지당 단일 객체 — 다중 객체 일반화 미확인", 1),
            ("2-class binary — 다른 병해와 구분 불가", 1),
            ("도메인 shift 미검증", 0, True),
            ("카메라/지역/시기 기반 split 실험 부재", 1),
            ("MPS 벤치마킹 특이사항", 0, True),
            ("일부 연산 CPU fallback, 절대 latency는 환경 의존", 1),
        ])

    add_content_slide(prs, "향후 계획",
        [
            ("단기 가능", 0, True),
            ("Distillation 하이퍼파라미터 sweep (α, T)", 1),
            ("양자화: static PTQ 또는 QAT로 속도 개선 시도", 1),
            ("중기", 0, True),
            ("도메인 일반화 — Leave-one-camera-out / location-out 평가", 1),
            ("멀티모달 — 외부 기상 시계열 결합 (AI Hub 외부 CSV 필요)", 1),
            ("장기 / 별도 프로젝트", 0, True),
            ("병변 단위 라벨링 (canker lesion-level segmentation)", 1),
            ("CoreML export + iOS/iPadOS 앱 통합", 1),
            ("현장 데이터 수집 (야외, 다양한 조명)", 1),
        ])

    # --- 10. Summary ---
    add_section_header(prs, "요약", "개발 성과 · 학습 · 공개")

    add_content_slide(prs, "개발 성과 한눈에",
        [
            ("구현 완료 Phase", 0, True),
            ("P0 Common / P1 Classification / P2 Detection / P3 Segmentation", 1),
            ("P1b 백본 비교 / P2b YOLO variants / P3b Seg 아키텍처", 1),
            ("P4 장기학습(A1·A2) / 양자화(G) / Distillation(H)", 1),
            ("성능 Highlights", 0, True),
            ("Classification Acc 98.83%, canker Recall 100%", 1),
            ("Detection mAP@0.5 0.994 (5ep) → 0.9945 (30ep)", 1),
            ("Segmentation mIoU 0.9836 (20ep DeepLabV3+)", 1),
            ("배포 Highlights", 0, True),
            ("ViT-Small INT8: 크기 −73%, 정확도 유지", 1),
            ("MobileNetV3 264 FPS (모바일 후보)", 1),
        ])

    add_content_slide(prs, "재현성 & 공개",
        [
            ("코드 품질", 0, True),
            ("93개 pytest (unit + integration)", 1),
            ("~70개 atomic commits", 1),
            ("TDD + 설계 스펙 + Phase별 구현 plan 문서", 1),
            ("재현성", 0, True),
            ("environment.yml로 env 재현", 1),
            ("YAML config + seed + timestamped output + config snapshot", 1),
            ("공개", 0, True),
            ("GitHub: github.com/YBNML/disease_01 (public)", 1),
            ("README에 알고리즘 원리 · 학습 로직 · 결과 해석 전부 수록", 1),
            ("Jupyter 분석 노트북 8종 (GitHub 바로 렌더링)", 1),
        ])

    # --- Final ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.enum.shapes import MSO_SHAPE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_PRIMARY
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), prs.slide_width - Inches(1.4), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "감사합니다"
    _set_font(r, 60, bold=True, color=COLOR_HEADER_FG)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = "Q & A"
    _set_font(r, 28, color=COLOR_HEADER_FG)
    p3 = tf.add_paragraph()
    r = p3.add_run(); r.text = "github.com/YBNML/disease_01"
    _set_font(r, 16, color=COLOR_HEADER_FG)

    return prs


def main():
    prs = build_presentation()
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"presentation written: {OUT_PATH}")
    print(f"  slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
